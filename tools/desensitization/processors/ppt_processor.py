"""
PowerPoint 处理器（.pptx）
基于 python-pptx，遍历每张幻灯片的文本框、表格、占位符，
在 run 级别做替换以保留格式。
"""
import logging
from typing import List

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from core.llm_client import Replacement
from processors.base import BaseProcessor

logger = logging.getLogger(__name__)


class PPTProcessor(BaseProcessor):
    """PowerPoint .pptx 处理器"""

    def extract_text(self, file_path: str) -> str:
        prs = Presentation(file_path)
        parts: List[str] = []

        for slide in prs.slides:
            for shape in slide.shapes:
                self._extract_from_shape(shape, parts)

        return "\n".join(parts)

    def _extract_from_shape(self, shape, parts: List[str]):
        """递归提取形状中的文本（含表格、组合形状）"""
        # 文本框
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs)
                if text:
                    parts.append(text)

        # 表格
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs)
                        if text:
                            parts.append(text)

        # 组合形状
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                self._extract_from_shape(child, parts)

    def apply_replacements(
        self, file_path: str, replacements: List[Replacement]
    ) -> List[Replacement]:
        prs = Presentation(file_path)
        applied_set = set()
        applied: List[Replacement] = []

        for slide in prs.slides:
            for shape in slide.shapes:
                self._replace_in_shape(shape, replacements, applied_set, applied)

        prs.save(file_path)
        return applied

    def _replace_in_shape(
        self,
        shape,
        replacements: List[Replacement],
        applied_set: set,
        applied: List[Replacement],
    ):
        """递归在形状中做替换"""
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                self._replace_in_paragraph(para, replacements, applied_set, applied)

        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        self._replace_in_paragraph(
                            para, replacements, applied_set, applied
                        )

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                self._replace_in_shape(child, replacements, applied_set, applied)
